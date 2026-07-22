"""PPT 解析器 — 基于 python-pptx 的真实实现"""
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from parser import ParserAdapter


class PptxParser(ParserAdapter):
    """PPTX 文档解析器

    使用 python-pptx 提取每页幻灯片的标题、正文文本和图片 OCR 文字。
    """

    def __init__(self, ocr=None):
        """初始化

        Args:
            ocr: OCRProvider 实例，为 None 时跳过图片识别
        """
        self.ocr = ocr

    def _extract_title(self, slide) -> str:
        """从 slide 中提取标题文本"""
        if slide.shapes.title and slide.shapes.title.text.strip():
            return slide.shapes.title.text.strip()

        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                return shape.text.strip().split("\n")[0][:100]
        return ""

    def _extract_content(self, slide) -> str:
        """从 slide 中提取正文文本（排除标题）"""
        texts = []
        title_text = self._extract_title(slide)

        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text and text != title_text:
                        texts.append(text)

        return "\n".join(texts)

    async def _extract_images_text(self, slide) -> str:
        """从 slide 的图片 shape 中提取 OCR 文字"""
        if self.ocr is None:
            return ""

        image_texts = []
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    image_data = shape.image.blob
                    ocr_text = await self.ocr.recognize(image_data)
                    if ocr_text.strip():
                        image_texts.append(ocr_text.strip())
                except (OSError, AttributeError):
                    # 单张图片 OCR 失败不阻塞整体解析
                    pass

        return "\n".join(image_texts)

    async def parse(self, file_path: str) -> list[dict]:
        """解析 PPTX 文件

        Returns:
            [{page, title, content}, ...]
        """
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        if path.suffix.lower() != ".pptx":
            raise ValueError(f"不支持的文件格式: {path.suffix}，需要 .pptx")

        prs = Presentation(str(path))
        pages = []

        for i, slide in enumerate(prs.slides):
            title = self._extract_title(slide)
            text_content = self._extract_content(slide)
            ocr_content = await self._extract_images_text(slide)

            # 合并文本和 OCR 结果
            content_parts = []
            if text_content:
                content_parts.append(text_content)
            if ocr_content:
                content_parts.append(f"[图片文字]\n{ocr_content}")

            pages.append({
                "page": i + 1,
                "title": title,
                "content": "\n".join(content_parts),
            })

        return pages

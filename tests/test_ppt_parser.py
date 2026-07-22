"""测试：PPT 解析器"""
import asyncio
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from parser import get_parser
from parser.ppt_parser import PptxParser


def _create_test_pptx(pages: int = 3) -> Path:
    """创建用于测试的 PPTX 文件"""
    prs = Presentation()
    for i in range(pages):
        slide_layout = prs.slide_layouts[0]  # title slide layout
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = f"第{i+1}章 标题"

        # 添加正文文本框
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(3)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = f"这是第{i+1}页的正文内容。"
        tf.add_paragraph().text = f"第二个段落 - 页码 {i+1}"

    tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
    prs.save(tmp.name)
    tmp.close()
    return Path(tmp.name)


class TestPptxParser:
    """PptxParser 真实解析测试"""

    def test_parse_returns_all_pages(self):
        pptx_path = _create_test_pptx(pages=3)
        parser = PptxParser()
        result = asyncio.run(parser.parse(str(pptx_path)))

        assert len(result) == 3
        for i, page in enumerate(result):
            assert page["page"] == i + 1
            assert isinstance(page["title"], str)
            assert isinstance(page["content"], str)
        pptx_path.unlink()

    def test_parse_extracts_title(self):
        pptx_path = _create_test_pptx(pages=1)
        parser = PptxParser()
        result = asyncio.run(parser.parse(str(pptx_path)))

        assert "第1章 标题" in result[0]["title"]
        pptx_path.unlink()

    def test_parse_extracts_content(self):
        pptx_path = _create_test_pptx(pages=1)
        parser = PptxParser()
        result = asyncio.run(parser.parse(str(pptx_path)))

        assert "这是第1页的正文内容" in result[0]["content"]
        pptx_path.unlink()

    def test_parse_empty_pptx_returns_empty_list(self):
        """无幻灯片的 PPTX 返回空列表"""
        prs = Presentation()
        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)
        tmp.close()

        parser = PptxParser()
        result = asyncio.run(parser.parse(tmp.name))
        assert result == []
        Path(tmp.name).unlink()

    def test_file_not_found_raises(self):
        parser = PptxParser()
        with pytest.raises(FileNotFoundError):
            asyncio.run(parser.parse("/nonexistent/path.pptx"))

    def test_wrong_extension_raises(self):
        tmp = tempfile.NamedTemporaryFile(suffix=".txt", delete=False)
        tmp.write(b"not a pptx")
        tmp.close()

        parser = PptxParser()
        with pytest.raises(ValueError, match="不支持的文件格式"):
            asyncio.run(parser.parse(tmp.name))
        Path(tmp.name).unlink()

    def test_multi_page_title_content_distinct(self):
        """不同页面的标题和内容应不同"""
        pptx_path = _create_test_pptx(pages=5)
        parser = PptxParser()
        result = asyncio.run(parser.parse(str(pptx_path)))

        titles = [p["title"] for p in result]
        # 每页标题应包含不同页码
        assert len(set(titles)) == len(titles)
        pptx_path.unlink()


class TestGetParserFactory:
    """get_parser 工厂函数"""

    def test_mock_mode_returns_mock_parser(self):
        import os
        os.environ["PARSER_MODE"] = "mock"
        parser = get_parser("test.pptx")
        assert parser.__class__.__name__ == "MockParser"
        os.environ.pop("PARSER_MODE")

    def test_pptx_mode_returns_pptx_parser(self):
        import os
        os.environ["PARSER_MODE"] = "python_pptx"
        parser = get_parser("test.pptx")
        assert parser.__class__.__name__ == "PptxParser"
        os.environ.pop("PARSER_MODE")

    def test_mock_parser_returns_structure(self):
        import os
        os.environ["PARSER_MODE"] = "mock"
        parser = get_parser()
        result = asyncio.run(parser.parse("dummy.pptx"))
        assert len(result) == 5
        assert all("page" in p and "title" in p and "content" in p for p in result)
        os.environ.pop("PARSER_MODE")

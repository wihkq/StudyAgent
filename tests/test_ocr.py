"""测试：OCR 图片文字识别"""
import asyncio
import tempfile
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from parser.ocr import MockOCR, OCRProvider, get_ocr
from parser.ppt_parser import PptxParser


class TestMockOCR:
    """MockOCR 基本行为"""

    def test_recognize_non_empty_returns_mock_text(self):
        ocr = MockOCR()
        result = asyncio.run(ocr.recognize(b"fake image data"))
        assert len(result) > 0
        assert "图片文字" in result

    def test_recognize_empty_returns_empty(self):
        ocr = MockOCR()
        result = asyncio.run(ocr.recognize(b""))
        assert result == ""

    def test_ocr_provider_raises_not_implemented(self):
        base = OCRProvider()
        with pytest.raises(NotImplementedError):
            asyncio.run(base.recognize(b"test"))


class TestGetOCR:
    """get_ocr 工厂"""

    def test_default_returns_mock_ocr(self):
        ocr = get_ocr()
        assert isinstance(ocr, MockOCR)


class TestPptxParserWithOCR:
    """PptxParser OCR 集成"""

    def _create_pptx_with_image(self) -> Path:
        """创建包含图片的 PPTX"""
        import io
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # blank layout
        slide = prs.slides.add_slide(slide_layout)

        # 添加文本
        left = Inches(1)
        top = Inches(0.5)
        txBox = slide.shapes.add_textbox(left, top, Inches(8), Inches(1.5))
        tf = txBox.text_frame
        tf.text = "带图片的页面"
        tf.add_paragraph().text = "这段正文不会被标题过滤排除"

        # 添加一个 1x1 像素的 PNG 作为测试图片
        png_data = (
            b"\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01"
            b"\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\x0cIDATx\x9cc\xf8\x0f"
            b"\x00\x00\x01\x01\x00\x05\x18\xd8N\x00\x00\x00\x00IEND\xaeB`\x82"
        )
        left = Inches(1)
        top = Inches(2)
        slide.shapes.add_picture(io.BytesIO(png_data), left, top)

        tmp = tempfile.NamedTemporaryFile(suffix=".pptx", delete=False)
        prs.save(tmp.name)
        tmp.close()
        return Path(tmp.name)

    def test_parse_with_ocr_appends_image_text(self):
        pptx_path = self._create_pptx_with_image()
        ocr = MockOCR()
        parser = PptxParser(ocr=ocr)
        result = asyncio.run(parser.parse(str(pptx_path)))

        assert len(result) == 1
        content = result[0]["content"]
        assert "图片文字" in content
        assert "这段正文不会被标题过滤排除" in content
        pptx_path.unlink()

    def test_parse_without_ocr_skips_images(self):
        """无 OCR 时跳过图片，不崩溃"""
        pptx_path = self._create_pptx_with_image()
        parser = PptxParser(ocr=None)
        result = asyncio.run(parser.parse(str(pptx_path)))

        assert len(result) == 1
        content = result[0]["content"]
        assert "图片文字" not in content
        assert "这段正文不会被标题过滤排除" in content
        pptx_path.unlink()

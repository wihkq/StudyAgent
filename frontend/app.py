"""StudyAgent 前端 — 智能期末复习助手"""
import sys
from pathlib import Path
sys.path.insert(0, str(Path(__file__).parent.parent))

import datetime
import random
import time
import streamlit as st

# ===== 页面配置 =====
st.set_page_config(
    page_title="StudyAgent · 智能复习助手",
    page_icon="📚",
    layout="wide",
    initial_sidebar_state="expanded",
)

# ===== 自定义样式 =====
st.markdown("""
<style>
    /* 侧边栏 */
    [data-testid="stSidebar"] {
        background: linear-gradient(180deg, #1a1a2e 0%, #16213e 100%);
    }
    [data-testid="stSidebar"] * { color: #e0e0e0 !important; }
    [data-testid="stSidebar"] .stRadio > div { gap: 6px; }
    [data-testid="stSidebar"] .stRadio label {
        padding: 10px 16px; border-radius: 10px; font-size: 15px;
        transition: all 0.2s;
    }
    [data-testid="stSidebar"] .stRadio label:hover {
        background: rgba(255,255,255,0.08);
    }

    /* 主标题 */
    h1 { color: #1a1a2e !important; font-weight: 700 !important; }

    /* 卡片区 */
    .card-row { display: flex; gap: 16px; margin-bottom: 24px; }
    .stat-card {
        flex: 1; padding: 24px 20px; border-radius: 16px;
        background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: #fff; text-align: center;
    }
    .stat-card.green  { background: linear-gradient(135deg, #11998e 0%, #38ef7d 100%); }
    .stat-card.orange { background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); }
    .stat-card.blue   { background: linear-gradient(135deg, #4facfe 0%, #00f2fe 100%); }
    .stat-card .value { font-size: 36px; font-weight: 800; }
    .stat-card .label { font-size: 14px; opacity: 0.85; margin-top: 4px; }

    /* 课程卡片 */
    .course-card {
        padding: 20px 24px; border-radius: 14px; background: #fff; margin-bottom: 12px;
        box-shadow: 0 2px 12px rgba(0,0,0,0.06); border-left: 5px solid #667eea;
        transition: transform 0.15s;
    }
    .course-card:hover { transform: translateY(-2px); box-shadow: 0 4px 20px rgba(0,0,0,0.12); }

    /* 按钮 */
    .stButton button {
        border-radius: 10px !important; font-weight: 600 !important;
        transition: all 0.2s !important;
    }

    /* 首页大标题 */
    .hero-title { font-size: 42px; font-weight: 800; background: linear-gradient(135deg, #667eea, #764ba2); -webkit-background-clip: text; -webkit-text-fill-color: transparent; }
    .hero-sub { font-size: 18px; color: #666; margin-bottom: 32px; }
</style>
""", unsafe_allow_html=True)

# ===== 初始化 session_state =====
if "courses" not in st.session_state:
    st.session_state.courses = [
        {
            "id": "crs-demo01",
            "name": "示例课程",
            "created_date": datetime.date.today().isoformat(),
            "exam_date": (datetime.date.today() + datetime.timedelta(days=14)).isoformat(),
            "files": [
                {"name": "第一章课件.pptx", "upload_date": datetime.date.today().isoformat(), "size_mb": 2.3},
                {"name": "第二章课件.pptx", "upload_date": datetime.date.today().isoformat(), "size_mb": 3.1},
                {"name": "复习资料.pdf", "upload_date": datetime.date.today().isoformat(), "size_mb": 1.5},
            ],
        }
    ]
if "current_course_id" not in st.session_state:
    st.session_state.current_course_id = None
if "global_store" not in st.session_state:
    from knowledge.embedding import get_embedding
    from knowledge.vector_store import get_vector_store
    st.session_state.global_embedding = get_embedding()
    st.session_state.global_store = get_vector_store(dim=128)
    st.session_state.knowledge_count = 0
if "learning_stats" not in st.session_state:
    st.session_state.learning_stats = {
        "total_study_hours": 6.5,
        "completed_chapters": 3,
        "total_chapters": 8,
        "correct_rate": 0.78,
        "questions_answered": 24,
    }
if "exam_date" not in st.session_state:
    st.session_state.exam_date = datetime.date.today() + datetime.timedelta(days=14)


# ===== 辅助函数 =====
def add_course(name: str) -> str:
    """创建新课程，返回课程 ID"""
    import uuid
    cid = f"crs-{uuid.uuid4().hex[:8]}"
    st.session_state.courses.insert(0, {
        "id": cid,
        "name": name,
        "created_date": datetime.date.today().isoformat(),
        "exam_date": (datetime.date.today() + datetime.timedelta(days=14)).isoformat(),
        "files": [],
    })
    return cid

def add_file_to_course(course_id: str, filename: str, filedata: bytes) -> None:
    """向课程添加课件文件，同时保存到磁盘"""
    import uuid
    from pathlib import Path
    ext = Path(filename).suffix.lower()
    file_id = f"doc-{uuid.uuid4().hex[:8]}"
    save_dir = Path("uploads")
    save_dir.mkdir(exist_ok=True)
    save_path = save_dir / f"{file_id}{ext}"
    save_path.write_bytes(filedata)
    size_mb = round(len(filedata) / (1024 * 1024), 2)
    for c in st.session_state.courses:
        if c["id"] == course_id:
            c["files"].append({
                "name": filename,
                "file_id": file_id,
                "path": str(save_path),
                "upload_date": datetime.date.today().isoformat(),
                "size_mb": size_mb,
            })
            break


# ===== 侧边栏 =====
with st.sidebar:
    st.markdown("""
    <div style="text-align:center;padding:16px 0 8px">
        <div style="font-size:40px">📚</div>
        <div style="font-size:20px;font-weight:700;color:#fff">StudyAgent</div>
        <div style="font-size:12px;color:#aaa">智能期末复习助手</div>
    </div>
    """, unsafe_allow_html=True)
    st.markdown("---")

    page = st.radio(
        "导航",
        ["📋 首页概览", "📚 我的课程", "💬 AI 问答", "📅 复习计划", "📝 模拟考试", "🗺️ 知识地图", "📊 学习进度", "📝 错题本"],
        label_visibility="collapsed",
    )

    st.markdown("---")
    days_left = (st.session_state.exam_date - datetime.date.today()).days
    if days_left >= 0:
        st.markdown(f"""
        <div style="text-align:center;padding:12px;border-radius:12px;background:rgba(255,255,255,0.08)">
            <div style="font-size:12px;color:#aaa">距离考试</div>
            <div style="font-size:28px;font-weight:700;color:#38ef7d">{days_left} 天</div>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.error("考试日期已过")
    import os
    llm_mode = os.getenv("LLM_MODE", "mock")
    emoji = {"mock": "🧪", "deepseek": "🧠", "kimi": "🧠"}.get(llm_mode, "🤖")
    st.caption(f"{emoji} LLM: {llm_mode} | 📚 知识库: {st.session_state.knowledge_count} 条")


# ===== 页面路由 =====

if page == "📋 首页概览":
    st.markdown('<div class="hero-title">StudyAgent</div>', unsafe_allow_html=True)
    st.markdown('<div class="hero-sub">基于多模态 RAG 的智能期末复习助手</div>', unsafe_allow_html=True)

    # 统计卡片
    stats = st.session_state.learning_stats
    total_files = sum(len(c.get("files", [])) for c in st.session_state.courses)
    st.markdown(f"""
    <div class="card-row">
        <div class="stat-card">
            <div class="value">{len(st.session_state.courses)}</div>
            <div class="label">已创建课程</div>
        </div>
        <div class="stat-card green">
            <div class="value">{stats['completed_chapters']}/{stats['total_chapters']}</div>
            <div class="label">已完成章节</div>
        </div>
        <div class="stat-card orange">
            <div class="value">{stats['correct_rate']:.0%}</div>
            <div class="label">答题正确率</div>
        </div>
        <div class="stat-card blue">
            <div class="value">{stats['total_study_hours']}h</div>
            <div class="label">累计学习时长</div>
        </div>
    </div>
    """, unsafe_allow_html=True)

    # 考试日期
    col_a, col_b = st.columns([1, 3])
    with col_a:
        new_date = st.date_input("📅 考试日期", value=st.session_state.exam_date)
        if new_date != st.session_state.exam_date:
            st.session_state.exam_date = new_date
            st.rerun()

    st.divider()

    # 最近课程
    st.subheader("📖 最近课程")
    if st.session_state.courses:
        for c in st.session_state.courses[:3]:
            fc = len(c.get("files", []))
            st.markdown(f"""
            <div class="course-card">
                <strong style="font-size:17px">📚 {c['name']}</strong><br>
                <span style="color:#888;font-size:13px">
                    {fc} 个课件 · 创建于 {c['created_date']}
                </span>
            </div>
            """, unsafe_allow_html=True)
    else:
        st.info("还没有课程，去「上传课件」开始吧！")


elif page == "📚 我的课程":
    st.title("📚 我的课程")
    st.caption(f"共 {len(st.session_state.courses)} 门课程")

    # 新建课程
    with st.expander("➕ 创建新课程"):
        c1, c2 = st.columns(2)
        with c1:
            new_name = st.text_input("课程名称", placeholder="例如：高等数学", key="new_course_name")
        with c2:
            new_exam = st.date_input("考试日期", value=datetime.date.today() + datetime.timedelta(days=14), key="new_exam_date")
        if st.button("创建课程", type="primary", use_container_width=True):
            if new_name.strip():
                cid = add_course(new_name.strip())
                st.session_state.courses[0]["exam_date"] = new_exam.isoformat()
                st.success(f"「{new_name}」创建成功！")
                st.rerun()
            else:
                st.error("请输入课程名称")

    if not st.session_state.courses:
        st.info("还没有课程，点击上方「创建新课程」开始吧！")
    else:
        # 课程列表 — 点击进入详情
        for course in st.session_state.courses:
            fc = len(course.get("files", []))
            with st.container(border=True):
                c1, c2, c3 = st.columns([3, 1, 1])
                with c1:
                    st.markdown(f"### 📚 {course['name']}")
                    st.caption(f"创建于 {course['created_date']} · 考试 {course.get('exam_date', '未设置')} · {fc} 个课件")
                with c2:
                    if st.button("📂 进入课程", key=f"enter_{course['id']}", use_container_width=True):
                        st.session_state.current_course_id = course["id"]
                        st.rerun()
                with c3:
                    if st.button("🗑️ 删除", key=f"del_{course['id']}", use_container_width=True):
                        st.session_state.courses = [c for c in st.session_state.courses if c["id"] != course["id"]]
                        if st.session_state.current_course_id == course["id"]:
                            st.session_state.current_course_id = None
                        st.rerun()

        # ===== 课程详情（选中课程时显示）=====
        if st.session_state.current_course_id:
            cur = next((c for c in st.session_state.courses if c["id"] == st.session_state.current_course_id), None)
            if cur:
                st.divider()
                st.subheader(f"📂 {cur['name']} — 课程详情")

                # 上传课件到当前课程
                uploaded_files = st.file_uploader(
                    f"上传课件到「{cur['name']}」",
                    type=["pptx", "pdf"],
                    accept_multiple_files=True,
                    key=f"upload_{cur['id']}",
                )
                if uploaded_files:
                    for f in uploaded_files:
                        add_file_to_course(cur["id"], f.name, f.read())
                    st.success(f"已添加 {len(uploaded_files)} 个课件到「{cur['name']}」")
                    # 自动构建知识库
                    try:
                        from knowledge.chunker import chunk_pages
                        from knowledge.retriever import Retriever
                        import asyncio, io
                        all_pages = []
                        for fdata in uploaded_files:
                            if fdata.name.endswith(".pptx"):
                                from pptx import Presentation
                                fdata.seek(0)
                                prs = Presentation(io.BytesIO(fdata.read()))
                                for i, slide in enumerate(prs.slides):
                                    title = slide.shapes.title.text if slide.shapes.title else ""
                                    texts = []
                                    for s in slide.shapes:
                                        if s.has_text_frame:
                                            for p in s.text_frame.paragraphs:
                                                if p.text.strip():
                                                    texts.append(p.text.strip())
                                    all_pages.append({"page": i+1, "title": title, "content": " ".join(texts)})
                        if all_pages:
                            chunks = chunk_pages(all_pages, source_file=fdata.name)
                            ret = Retriever(st.session_state.global_embedding, st.session_state.global_store)
                            asyncio.run(ret.add_chunks(chunks))
                            st.session_state.knowledge_count += len(chunks)
                            st.success(f"✅ 知识库已更新，新增 {len(chunks)} 条知识")
                    except Exception as e:
                        st.warning(f"知识库构建失败（问答仍可用 Mock 模式）: {e}")
                    st.rerun()

                # 已上传课件列表
                if cur["files"]:
                    st.markdown("**已上传课件：**")
                    for f in cur["files"]:
                        st.markdown(f"- 📄 {f['name']} ({f['size_mb']} MB) — {f['upload_date']}")
                else:
                    st.info("暂无课件，请上传。")

                # PPT 预览
                if "view_file" in st.session_state and st.session_state.view_file:
                    vf = st.session_state.view_file
                    if vf.get("path", "").endswith(".pptx"):
                        st.divider()
                        st.subheader(f"📖 预览：{vf['name']}")
                        try:
                            from pptx import Presentation
                            prs = Presentation(vf["path"])
                            for i, slide in enumerate(prs.slides):
                                with st.container(border=True):
                                    title = slide.shapes.title.text if slide.shapes.title else f"第{i+1}页"
                                    st.markdown(f"**{i+1}. {title}**")
                                    texts = []
                                    for shape in slide.shapes:
                                        if shape.has_text_frame and not (slide.shapes.title and shape == slide.shapes.title):
                                            for para in shape.text_frame.paragraphs:
                                                if para.text.strip():
                                                    texts.append(para.text.strip())
                                    if texts:
                                        st.caption(" | ".join(texts[:5]))
                        except Exception as e:
                            st.error(f"解析失败：{e}")
                    else:
                        st.info("仅支持预览 .pptx 文件，PDF 预览后续支持。")
                    if st.button("✕ 关闭预览", key="close_preview", use_container_width=True):
                        st.session_state.view_file = None
                        st.rerun()

                if st.button("✕ 关闭课程详情", use_container_width=True):
                    st.session_state.current_course_id = None
                    st.session_state.view_file = None
                    st.rerun()


elif page == "💬 AI 问答":
    st.title("💬 智能问答")
    st.caption("基于已上传课件内容，回答你的学习问题")

    if "chat_history" not in st.session_state:
        st.session_state.chat_history = []

    # 聊天记录
    for msg in st.session_state.chat_history[-10:]:
        with st.chat_message(msg["role"]):
            st.markdown(msg["content"])

    question = st.chat_input("输入你的问题，例如：什么是局部性原理？")
    if question:
        st.chat_message("user").markdown(question)
        try:
            from agents.tutor_agent import get_tutor_agent
            from knowledge.retriever import Retriever
            emb = st.session_state.global_embedding
            store = st.session_state.global_store
            retriever = Retriever(emb, store)
            import asyncio
            chunks = asyncio.run(retriever.retrieve(question, top_k=3))
            agent = get_tutor_agent()
            result = asyncio.run(agent.answer(question, chunks))
            answer = result["answer"]
            st.session_state.chat_history.append({"role": "user", "content": question})
            st.session_state.chat_history.append({"role": "assistant", "content": answer})
            st.rerun()
        except Exception as e:
            st.error(f"问答出错：{e}\n\n请确认已上传课件并构建知识库。")


elif page == "📅 复习计划":
    st.title("📅 复习计划")
    st.caption("根据考试日期和可用时间，自动生成个性化复习安排")

    c1, c2 = st.columns(2)
    with c1:
        exam_date = st.date_input("考试日期", value=st.session_state.exam_date)
    with c2:
        daily_hours = st.slider("每日学习时间（小时）", 0.5, 12.0, 4.0, 0.5)

    chapters_data = [
        {"title": "第一章", "importance": 3, "page_count": 2, "key_points": ["核心概念"]},
        {"title": "第二章", "importance": 5, "page_count": 3, "key_points": ["重点理论", "关键公式"]},
        {"title": "第三章", "importance": 4, "page_count": 2, "key_points": ["核心方法"]},
        {"title": "第四章", "importance": 5, "page_count": 3, "key_points": ["进阶应用", "案例分析"]},
    ]

    if st.button("生成复习计划", type="primary", use_container_width=True):
        try:
            from agents.planner_agent import get_planner_agent
            agent = get_planner_agent()
            plan = __import__("asyncio").run(agent.generate(exam_date.isoformat(), daily_hours, chapters_data))
            st.success(plan["summary"]["message"])
            for day in plan["plan"]:
                with st.container(border=True):
                    st.markdown(f"**第 {day['day']} 天** — {day['date']}")
                    for t in day["tasks"]:
                        st.markdown(f"- {t['type']}：{t['title']}（{t['hours']}h）")
        except Exception as e:
            st.error(f"生成失败：{e}")


elif page == "📝 模拟考试":
    st.title("📝 模拟考试")
    st.caption("基于课程内容自动生成试卷，答题后自动批改")

    chapters_data = [
        {"title": "第一章", "importance": 3, "page_count": 2, "key_points": ["核心概念的定义", "学习目标"]},
        {"title": "第二章", "importance": 5, "page_count": 3, "key_points": ["重点理论", "关键公式推导", "应用场景"]},
        {"title": "第三章", "importance": 4, "page_count": 2, "key_points": ["核心方法论"]},
        {"title": "第四章", "importance": 5, "page_count": 3, "key_points": ["进阶技巧", "案例分析"]},
    ]

    if "exam_questions" not in st.session_state:
        st.session_state.exam_questions = None
        st.session_state.exam_answers = None

    tab1, tab2 = st.tabs(["📝 生成试卷", "📊 批改结果"])

    with tab1:
        q_count = st.slider("题目数量", 2, 10, 4)
        if st.button("开始生成试卷", type="primary", use_container_width=True):
            try:
                from agents.examiner_agent import get_examiner_agent
                agent = get_examiner_agent()
                exam = __import__("asyncio").run(agent.generate_exam(chapters_data, q_count))
                st.session_state.exam_questions = exam["questions"]
                st.session_state.exam_answers = [""] * len(exam["questions"])
                st.rerun()
            except Exception as e:
                st.error(f"生成失败：{e}")

        if st.session_state.exam_questions:
            st.subheader("请作答")
            for i, q in enumerate(st.session_state.exam_questions):
                st.markdown(f"**{i+1}.** [{q['type']}] {q['question']}")
                if q["type"] == "choice":
                    st.session_state.exam_answers[i] = st.radio(
                        f"选择答案", q["options"], key=f"ans_{i}", index=None
                    )
                else:
                    st.session_state.exam_answers[i] = st.text_input(
                        f"你的回答", key=f"ans_{i}", placeholder="请输入答案..."
                    )
            if st.button("提交批改", type="primary", use_container_width=True):
                try:
                    from agents.examiner_agent import get_examiner_agent
                    agent = get_examiner_agent()
                    result = __import__("asyncio").run(
                        agent.grade(st.session_state.exam_questions, st.session_state.exam_answers)
                    )
                    st.session_state.grade_result = result
                    st.rerun()
                except Exception as e:
                    st.error(f"批改失败：{e}")

    with tab2:
        if "grade_result" in st.session_state and st.session_state.grade_result:
            r = st.session_state.grade_result
            st.metric("得分", f"{r['score']}/{r['total']}（{r['percentage']}%）")
            for item in r["results"]:
                icon = "✅" if item["is_correct"] else "❌"
                with st.expander(f"{icon} {item['question'][:40]}..."):
                    st.markdown(f"你的答案：{item['your_answer']}")
                    st.markdown(f"正确答案：{item['correct_answer']}")
        else:
            st.info("请先在「生成试卷」中答题并提交批改。")


elif page == "📚 我的课程":

    st.title("📚 我的课程")
    st.caption(f"共 {len(st.session_state.courses)} 门课程")

    if not st.session_state.courses:
        st.info("还没有课程，去「上传课件」添加吧！")
    else:
        for course in st.session_state.courses:
            s = {"ready": ("✅", "已就绪"), "processing": ("⏳", "处理中"), "error": ("❌", "出错")}
            icon, stext = s.get(course["status"], ("❓", course["status"]))
            st.markdown(f"""
            <div class="course-card">
                <table style="width:100%">
                    <tr>
                        <td style="width:60%">
                            <span style="font-size:18px;font-weight:600">{icon} {course['name']}</span><br>
                            <span style="color:#888;font-size:13px">{stext} · {course['file_count']} 个课件 · {course['chapter_count']} 个章节</span>
                        </td>
                        <td style="text-align:center;font-size:22px;font-weight:700">{course['file_count']}<br><span style="font-size:12px;color:#888">课件</span></td>
                        <td style="text-align:center;font-size:22px;font-weight:700">{course['chapter_count']}<br><span style="font-size:12px;color:#888">章节</span></td>
                        <td style="text-align:right;color:#888;font-size:13px">{course['upload_date']}</td>
                    </tr>
                </table>
            </div>
            """, unsafe_allow_html=True)

        st.divider()
        if st.button("🗑️ 清空所有课程"):
            st.session_state.courses = []
            st.rerun()


elif page == "🗺️ 知识地图":
    st.title("🗺️ 知识地图")
    st.caption("课程结构一目了然，⭐ 越多代表越重要")

    chapters = [
        {"t": "第一章", "i": 3, "p": [1, 2]},
        {"t": "第二章", "i": 5, "p": [3, 4, 5]},
        {"t": "第三章", "i": 4, "p": [6, 7]},
        {"t": "第四章", "i": 5, "p": [8, 9]},
    ]
    for ch in chapters:
        bar = "█" * ch["i"] + "░" * (5 - ch["i"])
        st.markdown(f"""
        <div class="course-card" style="border-left-color: {'#f5576c' if ch['i']>=5 else '#4facfe' if ch['i']>=4 else '#667eea'}">
            <strong>{'⭐' * ch['i']} {ch['t']}</strong>
            <span style="float:right;color:#888;font-size:13px">第 {', '.join(str(p) for p in ch['p'])} 页</span>
            <div style="margin-top:6px;font-family:monospace;color:#aaa">{bar} {ch['i']}/5</div>
        </div>
        """, unsafe_allow_html=True)
    st.caption("💡 上传真实课件后，系统将自动生成课程知识地图。")


elif page == "📊 学习进度":
    st.title("📊 学习进度")
    stats = st.session_state.learning_stats

    c1, c2, c3 = st.columns(3)
    with c1:
        pct = stats["completed_chapters"] / stats["total_chapters"]
        st.metric("📖 章节完成率", f"{pct:.0%}")
        st.progress(pct, text=f"{stats['completed_chapters']}/{stats['total_chapters']} 章")
    with c2:
        st.metric("🎯 答题正确率", f"{stats['correct_rate']:.0%}")
        st.progress(stats["correct_rate"], text=f"{stats['questions_answered']} 题已答")
    with c3:
        st.metric("⏱️ 累计学习", f"{stats['total_study_hours']} 小时")

    st.divider()
    st.subheader("📅 最近七天学习记录")
    today = datetime.date.today()
    cols = st.columns(7)
    day_names = ["周一", "周二", "周三", "周四", "周五", "周六", "周日"]
    random.seed(42)
    for i, col in enumerate(cols):
        day = today - datetime.timedelta(days=6 - i)
        h = round(random.uniform(0, 3.5), 1)
        bar = "🟦" * max(1, int(h))
        with col:
            st.markdown(f"<div style='text-align:center;font-size:12px;color:#888'>{day_names[day.weekday()]}</div>", unsafe_allow_html=True)
            st.markdown(f"<div style='text-align:center;font-size:20px;font-weight:700'>{h}h</div>", unsafe_allow_html=True)
            st.markdown(f"<div style='text-align:center'>{bar}</div>", unsafe_allow_html=True)

    st.divider()
    with st.expander("🔧 模拟数据（演示用）"):
        if st.button("模拟学习 1 小时", use_container_width=True):
            stats["total_study_hours"] += 1
            if stats["completed_chapters"] < stats["total_chapters"]:
                stats["completed_chapters"] += 1
            stats["questions_answered"] += 5
            stats["correct_rate"] = min(1.0, stats["correct_rate"] + 0.02)
            st.rerun()
        if st.button("重置全部数据", use_container_width=True):
            st.session_state.learning_stats = {"total_study_hours": 0.0, "completed_chapters": 0, "total_chapters": 8, "correct_rate": 0.0, "questions_answered": 0}
            st.rerun()


elif page == "📝 错题本":
    st.title("📝 错题本")
    st.caption("自动收集错题，精准定位薄弱知识点")

    if "error_log" not in st.session_state:
        st.session_state.error_log = [
            {"q": "二进制与十进制的转换方法？", "a": "除以2取余", "c": "除2取余，逆向排列", "w": "进制转换"},
            {"q": "指令流水线的五个阶段？", "a": "3个阶段", "c": "取指·译码·执行·访存·写回", "w": "指令流水线"},
            {"q": "Cache 的工作原理？", "a": "直接访问", "c": "利用局部性原理缓存频繁访问数据", "w": "Cache 机制"},
        ]

    errs = st.session_state.error_log
    if errs:
        from collections import Counter
        wc = Counter(e["w"] for e in errs)
        st.subheader("🔍 薄弱知识点分布")
        for pt, cnt in wc.most_common():
            pct = cnt / len(errs)
            st.markdown(f"""
            <div style="margin:8px 0">
                <strong>{pt}</strong> <span style="color:#888">{cnt} 题 ({pct:.0%})</span>
                <div style="height:8px;border-radius:4px;background:{'#f5576c' if pct>0.5 else '#4facfe'};width:{int(pct*100)}%"></div>
            </div>
            """, unsafe_allow_html=True)

        st.divider()
        st.subheader("📋 错题详情")
        for i, e in enumerate(errs):
            with st.expander(f"错题 {i+1}：{e['q']}"):
                st.markdown(f"❌ 你的答案：**{e['a']}**")
                st.markdown(f"✅ 正确答案：**{e['c']}**")
                st.caption(f"📌 薄弱点：{e['w']}")

        st.divider()
        st.subheader("💡 复习建议")
        for pt, _ in wc.most_common(2):
            st.warning(f"建议重点复习「**{pt}**」相关章节，该知识点错误率较高。")
    else:
        st.info("暂无错题记录。去参加模拟考试，系统将自动收集错题。")
